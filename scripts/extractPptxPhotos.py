"""
extractPptxPhotos.py
────────────────────────────────────────────────────────────────────────────
Step 1 of 2: Extract photos from NEW_ASSETS.pptx into seed_photos/ folder.
Run BEFORE the Node seed script.

Usage:
    python scripts/extractPptxPhotos.py NEW_ASSETS.pptx

Output:
    seed_photos/
        asset_0/   ← Edo State (slides 1 & 2)
        asset_1/   ← Taraba State (slide 3)
        asset_2/   ← Adamawa State (slide 4)
        asset_3/   ← Niger State (slide 5)
        asset_4/   ← Nasarawa State (slides 6, 7 & 8)
"""

import sys, os, hashlib
from pptx import Presentation

# Slide number (1-based) → asset index
SLIDE_MAP = {
    1: 0,  # Edo
    2: 0,  # Edo (same site)
    3: 1,  # Taraba
    4: 2,  # Adamawa
    5: 3,  # Niger
    6: 4,  # Nasarawa
    7: 4,  # Nasarawa (same site)
    8: 4,  # Nasarawa (same site)
}

def main():
    pptx_path = sys.argv[1] if len(sys.argv) > 1 else 'NEW_ASSETS.pptx'
    if not os.path.exists(pptx_path):
        print(f'ERROR: File not found: {pptx_path}')
        sys.exit(1)

    out_dir = os.path.join(os.path.dirname(pptx_path), 'seed_photos')
    for i in range(5):
        os.makedirs(os.path.join(out_dir, f'asset_{i}'), exist_ok=True)

    prs = Presentation(pptx_path)
    seen   = set()
    counts = [0] * 5

    for slide_num, slide in enumerate(prs.slides, start=1):
        asset_idx = SLIDE_MAP.get(slide_num)
        if asset_idx is None:
            continue

        for shape in slide.shapes:
            if shape.shape_type != 13:   # 13 = PICTURE
                continue
            img   = shape.image
            blob  = img.blob
            h     = hashlib.md5(blob).hexdigest()[:10]
            if h in seen:
                continue                  # skip duplicate images across slides
            seen.add(h)

            ext   = img.ext or 'jpg'
            fname = f'slide{slide_num:02d}_{counts[asset_idx]:02d}_{h}.{ext}'
            dest  = os.path.join(out_dir, f'asset_{asset_idx}', fname)
            with open(dest, 'wb') as f:
                f.write(blob)
            counts[asset_idx] += 1
            print(f'  [slide {slide_num} → asset {asset_idx}] {fname}')

    print()
    for i, c in enumerate(counts):
        print(f'  asset_{i}: {c} photo(s)')
    print(f'\nDone. Photos saved to: {out_dir}')

if __name__ == '__main__':
    main()
